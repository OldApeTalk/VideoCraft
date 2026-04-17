"""Core pipeline for PPT-to-Video conversion.

Handles PPTX parsing (python-pptx), slide PNG export (PowerPoint COM),
and orchestrates TTS / subtitle / compose steps.
UI layer (tools/ppt2video) calls these functions from worker threads.
"""

from __future__ import annotations

import os
import shutil
import subprocess
from typing import Callable

ProgressCallback = Callable[[int, int, str], None]  # (done, total, msg)


# ── Step 1: import PPTX ─────────────────────────────────────────────────

def import_pptx(src: str, workdir: str) -> str:
    """Copy user-selected .pptx into workdir as source.pptx."""
    os.makedirs(workdir, exist_ok=True)
    dst = os.path.join(workdir, "source.pptx")
    shutil.copy2(src, dst)
    return dst


# ── Step 2a: extract speaker notes ──────────────────────────────────────

def extract_notes(pptx_path: str, notes_dir: str) -> list[str]:
    """Read per-slide speaker notes via python-pptx, write page_NN.txt.

    Returns list of output file paths (one per slide that has notes).
    Slides with empty notes still produce a file (empty).
    """
    from pptx import Presentation

    os.makedirs(notes_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    paths: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf:
                text = tf.text.strip()
        out = os.path.join(notes_dir, f"page_{i:02d}.txt")
        with open(out, "w", encoding="utf-8") as f:
            f.write(text)
        paths.append(out)

    return paths


# ── Step 2b: export slides to PNG via PowerPoint COM ────────────────────

def export_slides_to_png(
    pptx_path: str,
    pages_dir: str,
    *,
    on_progress: ProgressCallback | None = None,
) -> list[str]:
    """Export each slide as page_NN.png using PowerPoint COM automation.

    Requires Microsoft PowerPoint installed on Windows.
    Raises RuntimeError with a user-friendly message if COM fails.
    """
    try:
        import pythoncom
        from win32com.client import Dispatch
    except ImportError:
        raise RuntimeError(
            "pywin32 is not installed. "
            "Run: pip install pywin32"
        )

    os.makedirs(pages_dir, exist_ok=True)
    abs_pptx = os.path.abspath(pptx_path)
    abs_pages = os.path.abspath(pages_dir)

    pythoncom.CoInitialize()
    app = None
    pres = None
    try:
        try:
            app = Dispatch("PowerPoint.Application")
        except Exception:
            raise RuntimeError(
                "Microsoft PowerPoint not detected. "
                "PPT2Video currently requires Office installed on Windows."
            )

        pres = app.Presentations.Open(abs_pptx, WithWindow=False)
        total = pres.Slides.Count
        paths: list[str] = []

        for i in range(1, total + 1):
            out_path = os.path.join(abs_pages, f"page_{i:02d}.png")
            pres.Slides(i).Export(out_path, "PNG")
            paths.append(out_path)
            if on_progress:
                on_progress(i, total, f"page_{i:02d}.png")

        return paths

    finally:
        if pres:
            try:
                pres.Close()
            except Exception:
                pass
        if app:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


# ── Step 2 combined ─────────────────────────────────────────────────────

def run_step2(
    pptx_path: str,
    workdir: str,
    *,
    on_progress: ProgressCallback | None = None,
) -> tuple[list[str], list[str]]:
    """Extract notes + export PNGs. Returns (page_paths, note_paths)."""
    pages_dir = os.path.join(workdir, "pages")
    notes_dir = os.path.join(workdir, "notes")

    note_paths = extract_notes(pptx_path, notes_dir)
    page_paths = export_slides_to_png(pptx_path, pages_dir,
                                       on_progress=on_progress)
    return page_paths, note_paths


# ── Step 3: TTS ─────────────────────────────────────────────────────────

def synthesize_all_notes(
    notes_dir: str,
    audio_dir: str,
    voice_id: str,
    *,
    provider: str = "fish_audio",
    should_cancel: Callable[[], bool] | None = None,
    on_progress: ProgressCallback | None = None,
) -> list[str]:
    """Synthesize each page_NN.txt to page_NN.mp3 via core.tts."""
    from core import tts as core_tts

    os.makedirs(audio_dir, exist_ok=True)
    note_files = sorted(
        f for f in os.listdir(notes_dir)
        if f.startswith("page_") and f.endswith(".txt")
    )
    total = len(note_files)
    paths: list[str] = []

    for i, fname in enumerate(note_files):
        if should_cancel and should_cancel():
            raise InterruptedError("Cancelled")

        text_path = os.path.join(notes_dir, fname)
        with open(text_path, "r", encoding="utf-8") as f:
            text = f.read().strip()

        stem = os.path.splitext(fname)[0]
        out_path = os.path.join(audio_dir, f"{stem}.mp3")

        if on_progress:
            on_progress(i + 1, total, stem)

        if not text:
            paths.append("")
            continue

        try:
            core_tts.synthesize_text(
                text, out_path,
                voice_id=voice_id,
                provider=provider,
                should_cancel=should_cancel,
            )
        except Exception as e:
            raise RuntimeError(
                f"{stem}: {e}\n"
                f"(text length: {len(text)} chars)"
            ) from e
        paths.append(out_path)

    return paths


# ── Step 4: subtitle generation (character-ratio timeline) ──────────────

def generate_all_subs(
    notes_dir: str,
    audio_dir: str,
    subs_dir: str,
    *,
    on_progress: ProgressCallback | None = None,
) -> list[str]:
    """Generate per-page SRT from notes text + audio duration."""
    from core.srt_from_text import generate_srt_from_text

    os.makedirs(subs_dir, exist_ok=True)
    note_files = sorted(
        f for f in os.listdir(notes_dir)
        if f.startswith("page_") and f.endswith(".txt")
    )
    total = len(note_files)
    paths: list[str] = []

    for i, fname in enumerate(note_files):
        stem = os.path.splitext(fname)[0]
        text_path = os.path.join(notes_dir, fname)
        audio_path = os.path.join(audio_dir, f"{stem}.mp3")
        srt_path = os.path.join(subs_dir, f"{stem}.srt")

        if on_progress:
            on_progress(i + 1, total, stem)

        with open(text_path, "r", encoding="utf-8") as f:
            text = f.read().strip()

        if not text or not os.path.isfile(audio_path):
            paths.append("")
            continue

        generate_srt_from_text(text, audio_path, srt_path)
        paths.append(srt_path)

    return paths


# ── Step 5: compose MP4 ─────────────────────────────────────────────────

def compose_all(
    pages_dir: str,
    audio_dir: str,
    subs_dir: str,
    output_path: str,
    *,
    width: int = 1920,
    height: int = 1080,
    on_progress: ProgressCallback | None = None,
) -> str:
    """Compose per-page (image + audio + subtitle) segments, then concat."""
    from core.video_compose import compose_chapter
    from core.video_concat import concat_videos

    import tempfile

    page_files = sorted(
        f for f in os.listdir(pages_dir)
        if f.startswith("page_") and f.endswith(".png")
    )
    total = len(page_files)
    tmp_segments: list[str] = []

    try:
        for i, fname in enumerate(page_files):
            stem = os.path.splitext(fname)[0]
            image = os.path.join(pages_dir, fname)
            audio = os.path.join(audio_dir, f"{stem}.mp3")
            srt = os.path.join(subs_dir, f"{stem}.srt")

            if not os.path.isfile(audio):
                continue

            if on_progress:
                on_progress(i + 1, total, stem)

            tmp = tempfile.NamedTemporaryFile(
                delete=False, suffix=".mp4",
                dir=os.path.dirname(output_path))
            tmp.close()
            tmp_segments.append(tmp.name)

            compose_chapter(
                image, audio,
                srt if os.path.isfile(srt) else None,
                tmp.name,
                width=width, height=height,
            )

        if not tmp_segments:
            raise RuntimeError("No segments produced — check audio files")

        concat_videos(tmp_segments, output_path)

    finally:
        for f in tmp_segments:
            try:
                os.unlink(f)
            except OSError:
                pass

    return output_path
